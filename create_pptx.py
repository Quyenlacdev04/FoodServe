# -*- coding: utf-8 -*-
"""
FoodServe - Professional PowerPoint Presentation Generator
Tạo slide báo cáo đồ án chuyên nghiệp cho hội đồng bảo vệ
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============================================================
# COLOR PALETTE - Dark Premium Theme
# ============================================================
CLR_BG_DARK      = RGBColor(0x0F, 0x0F, 0x14)   # Nền chính đen xanh
CLR_BG_CARD      = RGBColor(0x1A, 0x1A, 0x24)   # Card nền
CLR_BG_CARD2     = RGBColor(0x22, 0x22, 0x30)   # Card nền phụ
CLR_ACCENT       = RGBColor(0xFF, 0x6B, 0x00)   # Cam chính (Primary)
CLR_ACCENT_LIGHT = RGBColor(0xFF, 0x8A, 0x33)   # Cam sáng
CLR_ACCENT2      = RGBColor(0xFF, 0xB8, 0x00)   # Vàng amber
CLR_GREEN        = RGBColor(0x10, 0xB9, 0x81)   # Xanh lá
CLR_BLUE         = RGBColor(0x38, 0x8B, 0xFC)   # Xanh dương
CLR_PURPLE       = RGBColor(0xA7, 0x8B, 0xFA)   # Tím
CLR_RED          = RGBColor(0xEF, 0x44, 0x44)   # Đỏ
CLR_CYAN         = RGBColor(0x06, 0xB6, 0xD4)   # Cyan
CLR_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # Trắng
CLR_TEXT         = RGBColor(0xE5, 0xE5, 0xEA)   # Text chính
CLR_TEXT_DIM     = RGBColor(0x9C, 0x9C, 0xA8)   # Text mờ
CLR_BORDER       = RGBColor(0x2A, 0x2A, 0x3A)   # Viền

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_fill(slide, left, top, width, height, color, corner_radius=0):
    """Add rounded rectangle shape with solid fill"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Set corner radius
    if corner_radius > 0:
        shape.adjustments[0] = corner_radius
    else:
        shape.adjustments[0] = 0.04
    return shape


def add_circle(slide, left, top, size, color):
    """Add circle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_gradient_bar(slide, left, top, width, height, color1, color2):
    """Add a colored bar (simulated gradient with solid color)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=CLR_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI", line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    """Add formatted text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Set anchor
    tf.paragraphs[0].alignment = alignment
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    
    # Line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=14,
                       color=CLR_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                       font_name="Segoe UI", line_spacing=1.4, bullet=False):
    """Add multi-line text with optional bullets"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        display_text = line
        is_bold = bold
        text_color = color
        text_size = font_size
        
        # Support for formatting markers
        if line.startswith("##"):
            display_text = line[2:].strip()
            is_bold = True
            text_color = CLR_ACCENT
            text_size = font_size + 2
        elif line.startswith("#"):
            display_text = line[1:].strip()
            is_bold = True
            text_color = CLR_WHITE
            text_size = font_size + 1
        elif line.startswith(">>"):
            display_text = line[2:].strip()
            text_color = CLR_TEXT_DIM
            text_size = font_size - 1
        
        if bullet and not line.startswith("#"):
            display_text = "•  " + display_text
        
        p.text = display_text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color
        p.font.bold = is_bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        
        # Line spacing
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    
    return txBox


def add_icon_card(slide, left, top, width, height, icon, title, desc,
                  icon_color=CLR_ACCENT, bg_color=CLR_BG_CARD):
    """Add card with icon, title, and description"""
    # Card background
    card = add_shape_with_fill(slide, left, top, width, height, bg_color, 0.06)
    
    # Icon circle
    icon_size = Inches(0.55)
    circle = add_circle(slide, left + Inches(0.3), top + Inches(0.3), icon_size, icon_color)
    # Icon emoji text on circle
    add_text_box(slide, left + Inches(0.3), top + Inches(0.28), icon_size, icon_size,
                 icon, font_size=18, color=CLR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(1.0), width - Inches(0.6), Inches(0.4),
                 title, font_size=13, color=CLR_WHITE, bold=True)
    
    # Description
    add_text_box(slide, left + Inches(0.3), top + Inches(1.35), width - Inches(0.6), height - Inches(1.5),
                 desc, font_size=10, color=CLR_TEXT_DIM, line_spacing=1.5)


def add_section_number(slide, left, top, number, color=CLR_ACCENT):
    """Add large section number"""
    add_text_box(slide, left, top, Inches(1.2), Inches(1),
                 f"{number:02d}", font_size=60, color=color, bold=True,
                 font_name="Segoe UI Black")


def add_decorative_line(slide, left, top, width, color=CLR_ACCENT):
    """Add thin decorative line"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_decorative_dots(slide, positions, color=CLR_ACCENT):
    """Add decorative dot pattern"""
    for (x, y, size) in positions:
        circle = add_circle(slide, x, y, size, color)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        # Set transparency
        try:
            solidFill = circle.fill._fill
            alpha_elem = solidFill.makeelement(qn('a:alpha'), {'val': '30000'})
            for srgb in solidFill.iter():
                if srgb.tag.endswith('}srgbClr'):
                    srgb.append(alpha_elem)
                    break
        except:
            pass


def make_slide_header(slide, section_num, title, subtitle=""):
    """Create consistent slide header"""
    # Top accent bar
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), CLR_ACCENT, CLR_ACCENT2)
    
    # Section number
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(0.8), Inches(0.5),
                 f"{section_num:02d}", font_size=14, color=CLR_ACCENT, bold=True,
                 font_name="Segoe UI Black")
    
    # Separator dot
    add_text_box(slide, Inches(1.35), Inches(0.35), Inches(0.3), Inches(0.5),
                 "—", font_size=14, color=CLR_TEXT_DIM, bold=True)
    
    # Title
    add_text_box(slide, Inches(1.65), Inches(0.3), Inches(8), Inches(0.55),
                 title.upper(), font_size=15, color=CLR_WHITE, bold=True,
                 font_name="Segoe UI Semibold")
    
    if subtitle:
        add_text_box(slide, Inches(1.65), Inches(0.7), Inches(8), Inches(0.4),
                     subtitle, font_size=10, color=CLR_TEXT_DIM)
    
    # Right side decorative element - FoodServe branding
    add_text_box(slide, Inches(11.0), Inches(0.35), Inches(2), Inches(0.4),
                 "🍽️ FoodServe", font_size=11, color=CLR_ACCENT, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    
    # Bottom thin line separator under header
    add_gradient_bar(slide, Inches(0.8), Inches(1.1), Inches(11.7), Pt(1), CLR_BORDER, CLR_BORDER)


# ============================================================
# SLIDE CREATION FUNCTIONS
# ============================================================

def create_title_slide(prs):
    """Slide 1: Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, CLR_BG_DARK)
    
    # Top accent bar
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CLR_ACCENT, CLR_ACCENT2)
    
    # Decorative circles
    add_decorative_dots(slide, [
        (Inches(0.5), Inches(0.5), Inches(3.5)),
        (Inches(10), Inches(4.5), Inches(4)),
        (Inches(-1), Inches(5), Inches(2.5)),
    ], CLR_ACCENT)
    
    # Small label
    label_bg = add_shape_with_fill(slide, Inches(1.2), Inches(1.5), Inches(3), Inches(0.45), CLR_BG_CARD2, 0.15)
    add_text_box(slide, Inches(1.4), Inches(1.53), Inches(2.8), Inches(0.4),
                 "📋  BÁO CÁO ĐỒ ÁN TỐT NGHIỆP", font_size=10, color=CLR_ACCENT, bold=True,
                 font_name="Segoe UI Semibold")
    
    # Main Title
    add_text_box(slide, Inches(1.2), Inches(2.3), Inches(7), Inches(1.2),
                 "FOODSERVE", font_size=64, color=CLR_WHITE, bold=True,
                 font_name="Segoe UI Black")
    
    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(3.5), Inches(8), Inches(0.8),
                 "Ứng dụng đặt & giao đồ ăn trực tuyến", font_size=22,
                 color=CLR_ACCENT_LIGHT, bold=False, font_name="Segoe UI")
    
    # Description
    add_text_box(slide, Inches(1.2), Inches(4.3), Inches(8), Inches(0.7),
                 "Nền tảng web fullstack hiện đại với React.js, Node.js & MongoDB\nTích hợp thanh toán, realtime tracking, AI chatbot & gamification",
                 font_size=13, color=CLR_TEXT_DIM, line_spacing=1.6)
    
    # Accent line
    add_decorative_line(slide, Inches(1.2), Inches(5.3), Inches(2), CLR_ACCENT)
    
    # Student info cards
    info_y = Inches(5.6)
    # GVHD
    add_text_box(slide, Inches(1.2), info_y, Inches(3.5), Inches(0.35),
                 "GIẢNG VIÊN HƯỚNG DẪN", font_size=9, color=CLR_TEXT_DIM, bold=True)
    add_text_box(slide, Inches(1.2), info_y + Inches(0.3), Inches(3.5), Inches(0.35),
                 "ThS. [Tên Giảng Viên]", font_size=14, color=CLR_WHITE, bold=True)
    
    # SVTH
    add_text_box(slide, Inches(5.0), info_y, Inches(3.5), Inches(0.35),
                 "SINH VIÊN THỰC HIỆN", font_size=9, color=CLR_TEXT_DIM, bold=True)
    add_text_box(slide, Inches(5.0), info_y + Inches(0.3), Inches(3.5), Inches(0.35),
                 "Vũ Văn Quyền", font_size=14, color=CLR_WHITE, bold=True)
    
    # Class / Year
    add_text_box(slide, Inches(8.5), info_y, Inches(3.5), Inches(0.35),
                 "NIÊN KHÓA", font_size=9, color=CLR_TEXT_DIM, bold=True)
    add_text_box(slide, Inches(8.5), info_y + Inches(0.3), Inches(3.5), Inches(0.35),
                 "2024 — 2026", font_size=14, color=CLR_WHITE, bold=True)
    
    # Bottom bar
    add_gradient_bar(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), CLR_ACCENT, CLR_ACCENT2)
    
    # Right side large emoji icon
    add_text_box(slide, Inches(9.5), Inches(1.8), Inches(3), Inches(3),
                 "🍔", font_size=140, alignment=PP_ALIGN.CENTER)


def create_agenda_slide(prs):
    """Slide 2: Agenda / Nội dung trình bày"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 0, "NỘI DUNG TRÌNH BÀY", "Cấu trúc bài báo cáo đồ án")
    
    items = [
        ("01", "Tổng quan đề tài", "Lý do, mục tiêu, phạm vi & phương pháp nghiên cứu", CLR_ACCENT),
        ("02", "Cơ sở lý thuyết", "Ngôn ngữ, công cụ & môi trường phát triển", CLR_BLUE),
        ("03", "Phân tích hệ thống", "Biểu đồ Usecase & Lược đồ quan hệ ER", CLR_GREEN),
        ("04", "Demo sản phẩm", "Trình diễn các chức năng chính của FoodServe", CLR_PURPLE),
        ("05", "Hạn chế & Định hướng", "Đánh giá hạn chế và hướng phát triển tương lai", CLR_CYAN),
        ("06", "Kết luận", "Tổng kết & lời cảm ơn", CLR_ACCENT2),
    ]
    
    start_y = Inches(1.6)
    col1_x = Inches(1.0)
    col2_x = Inches(7.0)
    
    for i, (num, title, desc, color) in enumerate(items):
        col = col1_x if i < 3 else col2_x
        row = i % 3
        y = start_y + row * Inches(1.75)
        
        # Card
        card = add_shape_with_fill(slide, col, y, Inches(5.5), Inches(1.4), CLR_BG_CARD, 0.06)
        
        # Number
        num_bg = add_shape_with_fill(slide, col + Inches(0.25), y + Inches(0.3), Inches(0.7), Inches(0.7), color, 0.12)
        add_text_box(slide, col + Inches(0.25), y + Inches(0.35), Inches(0.7), Inches(0.7),
                     num, font_size=22, color=CLR_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Segoe UI Black")
        
        # Title
        add_text_box(slide, col + Inches(1.2), y + Inches(0.25), Inches(3.8), Inches(0.4),
                     title, font_size=15, color=CLR_WHITE, bold=True)
        
        # Description
        add_text_box(slide, col + Inches(1.2), y + Inches(0.7), Inches(3.8), Inches(0.5),
                     desc, font_size=10, color=CLR_TEXT_DIM)


def create_overview_reason_slide(prs):
    """Slide 3: Lý do chọn đề tài"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 1, "LÝ DO CHỌN ĐỀ TÀI", "Tổng quan về đề tài")
    
    # Left side - Main reasons
    reasons = [
        ("📈", "Nhu cầu thị trường", "Thị trường giao đồ ăn trực tuyến tại Việt Nam\ntăng trưởng mạnh mẽ (GrabFood, ShopeeFood,\nBaemin). Năm 2025 đạt ~3 tỷ USD doanh thu.", CLR_ACCENT),
        ("⚡", "Chuyển đổi số F&B", "Ngành F&B đang chuyển đổi số mạnh mẽ.\nNhà hàng nhỏ cần nền tảng trực tuyến\nđể cạnh tranh và tiếp cận khách hàng.", CLR_BLUE),
        ("🎓", "Ứng dụng kiến thức", "Áp dụng toàn diện kiến thức Fullstack:\nReact.js, Node.js, MongoDB, Socket.IO,\nRESTful API, Payment Gateway integration.", CLR_GREEN),
        ("🚀", "Công nghệ hiện đại", "Xây dựng ứng dụng web hoàn chỉnh với\ncác công nghệ hot nhất: Realtime, AI Chatbot,\nGamification, Responsive Design.", CLR_PURPLE),
    ]
    
    for i, (icon, title, desc, color) in enumerate(reasons):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.1)
        y = Inches(1.5) + row * Inches(2.85)
        
        add_icon_card(slide, x, y, Inches(5.7), Inches(2.5), icon, title, desc, color)


def create_overview_objectives_slide(prs):
    """Slide 4: Mục tiêu đề tài"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 1, "MỤC TIÊU ĐỀ TÀI", "Tổng quan về đề tài")
    
    # Main objective
    main_card = add_shape_with_fill(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(1.2), Inches(1.55), Inches(0.5), Inches(0.5),
                 "🎯", font_size=28)
    add_text_box(slide, Inches(1.8), Inches(1.6), Inches(10), Inches(0.4),
                 "Mục tiêu tổng quát", font_size=16, color=CLR_ACCENT, bold=True)
    add_text_box(slide, Inches(1.8), Inches(2.0), Inches(10), Inches(0.5),
                 "Xây dựng nền tảng web đặt và giao đồ ăn trực tuyến FoodServe hoàn chỉnh, hiện đại, tích hợp đa tính năng phục vụ 4 nhóm người dùng: Khách hàng, Nhà hàng, Tài xế và Quản trị viên.",
                 font_size=12, color=CLR_TEXT, line_spacing=1.5)
    
    # Specific objectives
    objectives = [
        ("👤", "Khách hàng", "Đặt món, thanh toán online,\ntheo dõi realtime, đánh giá,\nvoucher & tích điểm", CLR_ACCENT),
        ("🍳", "Nhà hàng", "Quản lý menu, đơn hàng,\ndoanh thu, phản hồi đánh giá,\nphí duy trì & thống kê", CLR_GREEN),
        ("🛵", "Tài xế", "Nhận đơn, tracking GPS,\nđiều hướng Leaflet Map,\nthống kê thu nhập", CLR_BLUE),
        ("👑", "Quản trị", "Dashboard tổng quan,\nquản lý user/đơn hàng,\nvoucher, doanh thu, cài đặt", CLR_PURPLE),
    ]
    
    y = Inches(3.1)
    for i, (icon, title, desc, color) in enumerate(objectives):
        x = Inches(0.8) + i * Inches(3.05)
        add_icon_card(slide, x, y, Inches(2.8), Inches(2.8), icon, title, desc, color)
    
    # Extra features
    extras_bg = add_shape_with_fill(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.95), CLR_BG_CARD2, 0.04)
    add_text_box(slide, Inches(1.2), Inches(6.25), Inches(11), Inches(0.7),
                 "✨ Tính năng nổi bật: AI Chatbot hỗ trợ  •  Đặt nhóm Split Bill  •  Gamification (Vòng quay, Scratch Card)  •  Đăng ký gói ăn  •  Thanh toán VNPay/MoMo/ZaloPay  •  Chế độ ăn lành mạnh",
                 font_size=11, color=CLR_TEXT_DIM, line_spacing=1.5)


def create_overview_scope_slide(prs):
    """Slide 5: Đối tượng & phạm vi nghiên cứu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 1, "ĐỐI TƯỢNG & PHẠM VI NGHIÊN CỨU", "Tổng quan về đề tài")
    
    # Left column - Đối tượng
    add_shape_with_fill(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(1.2), Inches(1.65), Inches(4.5), Inches(0.4),
                 "🔍  ĐỐI TƯỢNG NGHIÊN CỨU", font_size=14, color=CLR_ACCENT, bold=True)
    add_decorative_line(slide, Inches(1.2), Inches(2.15), Inches(1.5), CLR_ACCENT)
    
    subjects = [
        "Hệ thống đặt và giao đồ ăn trực tuyến",
        "Quy trình quản lý đơn hàng từ đặt đến giao",
        "Các mô hình thanh toán trực tuyến (VNPay, MoMo, ZaloPay, PayOS)",
        "Công nghệ web hiện đại: React.js, Node.js, MongoDB",
        "Giao tiếp thời gian thực với Socket.IO",
        "Trải nghiệm người dùng (UX/UI) trên ứng dụng web",
        "Tích hợp AI Chatbot hỗ trợ khách hàng (Gemini API)",
        "Gamification trong ứng dụng thương mại điện tử",
    ]
    add_multiline_text(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(4),
                       subjects, font_size=11, color=CLR_TEXT, bullet=True, line_spacing=1.7)
    
    # Right column - Phạm vi
    add_shape_with_fill(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(7.2), Inches(1.65), Inches(4.5), Inches(0.4),
                 "📐  PHẠM VI NGHIÊN CỨU", font_size=14, color=CLR_BLUE, bold=True)
    add_decorative_line(slide, Inches(7.2), Inches(2.15), Inches(1.5), CLR_BLUE)
    
    scopes = [
        "Nền tảng: Ứng dụng web (Web Application)",
        "Kiến trúc: Client-Server (SPA + REST API)",
        "Frontend: React.js 18 + Redux Toolkit + Vite",
        "Backend: Node.js + Express.js",
        "Database: MongoDB Atlas (NoSQL)",
        "Realtime: Socket.IO (WebSocket)",
        "Thanh toán: VNPay, MoMo, ZaloPay, PayOS, Xu",
        "Bản đồ: Leaflet + OpenStreetMap",
        "Deploy: Vercel (FE) + Render (BE)",
    ]
    add_multiline_text(slide, Inches(7.2), Inches(2.4), Inches(4.8), Inches(4),
                       scopes, font_size=11, color=CLR_TEXT, bullet=True, line_spacing=1.65)


def create_overview_methodology_slide(prs):
    """Slide 6: Phương pháp nghiên cứu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 1, "PHƯƠNG PHÁP NGHIÊN CỨU", "Tổng quan về đề tài")
    
    methods = [
        ("📚", "Nghiên cứu tài liệu", "Tham khảo tài liệu chuyên ngành,\ndocumentation chính thức của React.js,\nNode.js, MongoDB, Socket.IO. Nghiên cứu\ncác hệ thống tương tự: GrabFood, ShopeeFood.", CLR_ACCENT),
        ("🔬", "Phân tích & thiết kế", "Phân tích yêu cầu hệ thống,\nthiết kế cơ sở dữ liệu NoSQL,\nvẽ biểu đồ Use Case, lược đồ ER.\nThiết kế giao diện UI/UX hiện đại.", CLR_BLUE),
        ("💻", "Thực nghiệm & triển khai", "Lập trình theo mô hình Agile,\nchia sprint phát triển từng module.\nTest trên nhiều thiết bị và trình duyệt.\nDeploy production trên Vercel & Render.", CLR_GREEN),
        ("✅", "Kiểm thử & đánh giá", "Kiểm thử chức năng end-to-end,\nkiểm thử giao diện responsive,\nkiểm thử hiệu suất API,\nđánh giá trải nghiệm người dùng thực.", CLR_PURPLE),
    ]
    
    for i, (icon, title, desc, color) in enumerate(methods):
        x = Inches(0.8) + i * Inches(3.05)
        y = Inches(1.5)
        add_icon_card(slide, x, y, Inches(2.8), Inches(3.5), icon, title, desc, color)
    
    # Process flow
    flow_y = Inches(5.4)
    add_shape_with_fill(slide, Inches(0.8), flow_y, Inches(11.7), Inches(1.3), CLR_BG_CARD, 0.04)
    
    steps = ["Thu thập\nyêu cầu", "Phân tích\n& Thiết kế", "Lập trình\nphát triển", "Kiểm thử\n& Debug", "Triển khai\nDeploy"]
    
    for i, step in enumerate(steps):
        sx = Inches(1.2) + i * Inches(2.3)
        # Step circle
        circle = add_shape_with_fill(slide, sx, flow_y + Inches(0.2), Inches(0.45), Inches(0.45), CLR_ACCENT if i < 5 else CLR_TEXT_DIM, 0.5)
        add_text_box(slide, sx, flow_y + Inches(0.22), Inches(0.45), Inches(0.45),
                     str(i + 1), font_size=14, color=CLR_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        # Step text
        add_text_box(slide, sx + Inches(0.55), flow_y + Inches(0.15), Inches(1.5), Inches(0.9),
                     step, font_size=10, color=CLR_TEXT, line_spacing=1.4)
        
        # Arrow
        if i < 4:
            add_text_box(slide, sx + Inches(1.95), flow_y + Inches(0.25), Inches(0.3), Inches(0.3),
                         "→", font_size=16, color=CLR_ACCENT, bold=True)


def create_theory_languages_slide(prs):
    """Slide 7: Ngôn ngữ lập trình sử dụng"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 2, "NGÔN NGỮ LẬP TRÌNH", "Cơ sở lý thuyết")
    
    langs = [
        ("HTML5", "Ngôn ngữ đánh dấu siêu\nvăn bản, xây dựng cấu trúc\nnội dung trang web.", "⚙️", CLR_ACCENT),
        ("CSS3", "Ngôn ngữ tạo kiểu cho\ntrang web, responsive design,\nanimation & flexbox/grid.", "🎨", CLR_BLUE),
        ("JavaScript", "Ngôn ngữ lập trình chính\ncho cả Frontend (React.js)\nvà Backend (Node.js).", "💛", CLR_ACCENT2),
        ("JSX", "Cú pháp mở rộng của JS\ntrong React.js, kết hợp\nHTML trong JavaScript.", "⚛️", CLR_CYAN),
    ]
    
    for i, (name, desc, icon, color) in enumerate(langs):
        x = Inches(0.8) + i * Inches(3.05)
        y = Inches(1.5)
        
        card = add_shape_with_fill(slide, x, y, Inches(2.8), Inches(2.4), CLR_BG_CARD, 0.06)
        
        # Icon
        add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5),
                     icon, font_size=24)
        
        # Name
        add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(2.2), Inches(0.4),
                     name, font_size=16, color=color, bold=True)
        
        # Description
        add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(2.2), Inches(1),
                     desc, font_size=10, color=CLR_TEXT_DIM, line_spacing=1.6)
    
    # Frameworks & Libraries section
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
                 "📦  THƯ VIỆN & FRAMEWORK CHÍNH", font_size=13, color=CLR_WHITE, bold=True)
    add_decorative_line(slide, Inches(0.8), Inches(4.6), Inches(2), CLR_ACCENT)
    
    libs = [
        ("React.js 18", "Frontend SPA Framework", CLR_CYAN),
        ("Redux Toolkit", "State Management", CLR_PURPLE),
        ("Framer Motion", "Animation Library", CLR_ACCENT),
        ("React Router v6", "Client-side Routing", CLR_RED),
        ("Node.js", "Backend Runtime", CLR_GREEN),
        ("Express.js", "REST API Framework", CLR_TEXT_DIM),
        ("Mongoose", "MongoDB ODM", CLR_GREEN),
        ("Socket.IO", "WebSocket Realtime", CLR_BLUE),
        ("TailwindCSS", "Utility-first CSS", CLR_CYAN),
        ("Recharts", "Data Visualization", CLR_ACCENT2),
    ]
    
    for i, (name, desc, color) in enumerate(libs):
        col = i % 5
        row = i // 5
        x = Inches(0.8) + col * Inches(2.42)
        y = Inches(4.85) + row * Inches(1.1)
        
        card = add_shape_with_fill(slide, x, y, Inches(2.2), Inches(0.85), CLR_BG_CARD2, 0.08)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(1.9), Inches(0.35),
                     name, font_size=11, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.45), Inches(1.9), Inches(0.35),
                     desc, font_size=8, color=CLR_TEXT_DIM)


def create_theory_tools_slide(prs):
    """Slide 8: Công cụ hỗ trợ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 2, "CÔNG CỤ HỖ TRỢ PHÁT TRIỂN", "Cơ sở lý thuyết")
    
    tools = [
        ("💻", "Visual Studio Code", "Code editor chính, hỗ trợ IntelliSense,\nGit integration, terminal tích hợp,\nhệ sinh thái extension phong phú.", CLR_BLUE),
        ("🌿", "Git & GitHub", "Hệ thống quản lý phiên bản (VCS),\nlưu trữ source code, collaboration,\nversion tracking & branching.", CLR_GREEN),
        ("📮", "Postman", "Công cụ test & debug REST API,\ntạo collection request, automated testing,\ndocumentation API endpoints.", CLR_ACCENT),
        ("🧭", "Chrome DevTools", "Debug frontend, inspect elements,\nnetwork monitoring, performance profiling,\nresponsive design testing.", CLR_ACCENT2),
        ("🗄️", "MongoDB Compass", "GUI quản lý MongoDB database,\nvisual query builder, schema analysis,\ndata import/export & monitoring.", CLR_GREEN),
        ("🎨", "Figma", "Thiết kế giao diện UI/UX,\nprototyping, wireframing,\ndesign system & components.", CLR_PURPLE),
    ]
    
    for i, (icon, title, desc, color) in enumerate(tools):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.05)
        y = Inches(1.5) + row * Inches(2.9)
        add_icon_card(slide, x, y, Inches(3.75), Inches(2.55), icon, title, desc, color)


def create_theory_environment_slide(prs):
    """Slide 9: Môi trường phát triển & triển khai"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 2, "MÔI TRƯỜNG PHÁT TRIỂN & TRIỂN KHAI", "Cơ sở lý thuyết")
    
    # Architecture diagram (text-based)
    arch_bg = add_shape_with_fill(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.5), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(1.2), Inches(1.5), Inches(5), Inches(0.4),
                 "🏗️  KIẾN TRÚC HỆ THỐNG", font_size=13, color=CLR_ACCENT, bold=True)
    
    # Architecture layers
    layers = [
        ("CLIENT", "React.js 18\nVite (Dev Server)\nTailwindCSS\nRedux Toolkit", CLR_BLUE, Inches(1.0)),
        ("API GATEWAY", "Express.js\nREST API\nSocket.IO\nCORS / Auth", CLR_ACCENT, Inches(3.8)),
        ("SERVICES", "Node.js Runtime\nGemini AI API\nPayOS / VNPay\nNodemailer", CLR_GREEN, Inches(6.6)),
        ("DATABASE", "MongoDB Atlas\nMongoose ODM\n16 Collections\nIndexing", CLR_PURPLE, Inches(9.4)),
    ]
    
    for (label, desc, color, x) in layers:
        card = add_shape_with_fill(slide, x, Inches(2.0), Inches(2.4), Inches(1.7), color, 0.06)
        # Make card semi-transparent appearance by using darker version
        add_text_box(slide, x + Inches(0.15), Inches(2.05), Inches(2.1), Inches(0.3),
                     label, font_size=9, color=CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(2.35), Inches(2.1), Inches(1.2),
                     desc, font_size=9, color=CLR_WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.4)
        
        # Arrow between
        if x < Inches(9):
            add_text_box(slide, x + Inches(2.35), Inches(2.55), Inches(0.5), Inches(0.3),
                         "→", font_size=18, color=CLR_TEXT_DIM, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Development environment
    envs = [
        ("🖥️", "Phát triển (Development)", "Node.js v20+, npm, Vite Dev Server\n(HMR - Hot Module Replacement)\nPort 5173 (FE) + Port 5000 (BE)", CLR_BLUE),
        ("☁️", "Triển khai (Production)", "Frontend: Vercel (Static Hosting + CDN)\nBackend: Render (Node.js Cloud)\nDatabase: MongoDB Atlas (Cloud)", CLR_GREEN),
        ("🔧", "Môi trường khác", "OS: Windows 10/11\nBrowser: Chrome, Firefox, Edge\nMobile: Responsive Web (PWA-ready)", CLR_PURPLE),
    ]
    
    for i, (icon, title, desc, color) in enumerate(envs):
        x = Inches(0.8) + i * Inches(4.05)
        y = Inches(4.2)
        add_icon_card(slide, x, y, Inches(3.75), Inches(2.7), icon, title, desc, color)


def create_analysis_usecase_slide(prs):
    """Slide 10: Biểu đồ Use Case"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 3, "BIỂU ĐỒ USE CASE", "Phân tích hệ thống")
    
    # 4 Actor columns
    actors = [
        ("👤", "KHÁCH HÀNG", [
            "Đăng ký / Đăng nhập",
            "Tìm kiếm nhà hàng & món ăn",
            "Xem thực đơn & dinh dưỡng",
            "Đặt món & thanh toán",
            "Theo dõi đơn hàng realtime",
            "Đặt nhóm (Split Bill)",
            "Đăng ký gói ăn định kỳ",
            "Đánh giá & nhận xét",
            "Chơi game săn Xu",
            "Sử dụng AI Chatbot",
            "Quản lý voucher & xu",
            "Xem lịch sử đơn hàng",
        ], CLR_ACCENT),
        ("🍳", "NHÀ HÀNG", [
            "Đăng ký đối tác",
            "Quản lý thực đơn (CRUD)",
            "Xác nhận / từ chối đơn",
            "Cập nhật trạng thái món",
            "Xem thống kê doanh thu",
            "Phản hồi đánh giá",
            "Thanh toán phí duy trì",
            "Quản lý ảnh bìa cửa hàng",
        ], CLR_GREEN),
        ("🛵", "TÀI XẾ", [
            "Đăng ký tài xế",
            "Bật/Tắt trạng thái online",
            "Nhận đơn giao hàng",
            "Cập nhật vị trí GPS",
            "Xác nhận giao hàng",
            "Xem thống kê thu nhập",
            "Quản lý hồ sơ tài xế",
        ], CLR_BLUE),
        ("👑", "QUẢN TRỊ VIÊN", [
            "Dashboard tổng quan",
            "Quản lý người dùng",
            "Duyệt đối tác / tài xế",
            "Quản lý đơn hàng",
            "Quản lý mã giảm giá",
            "Cấu hình hệ thống",
            "Quản lý phí duy trì",
            "Bật/Tắt bảo trì",
            "Xem analytics & reports",
        ], CLR_PURPLE),
    ]
    
    for i, (icon, title, usecases, color) in enumerate(actors):
        x = Inches(0.6) + i * Inches(3.1)
        y = Inches(1.4)
        w = Inches(2.9)
        
        # Actor card
        card = add_shape_with_fill(slide, x, y, w, Inches(5.5), CLR_BG_CARD, 0.04)
        
        # Actor icon & name
        actor_bg = add_shape_with_fill(slide, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.65), color, 0.08)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.55),
                     f"{icon}  {title}", font_size=12, color=CLR_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        
        # Use cases
        for j, uc in enumerate(usecases):
            uy = y + Inches(1.05) + j * Inches(0.35)
            # Small bullet
            add_text_box(slide, x + Inches(0.2), uy, w - Inches(0.4), Inches(0.3),
                         f"▸ {uc}", font_size=9, color=CLR_TEXT, line_spacing=1.2)


def create_analysis_er_slide(prs):
    """Slide 11: Lược đồ ER"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 3, "LƯỢC ĐỒ QUAN HỆ THỰC THỂ (ER)", "Phân tích hệ thống")
    
    # Entity cards - arranged in grid
    entities = [
        ("User", "name, email, password,\nphone, role, coins, spins,\ntotalSpent, avatar", CLR_ACCENT, "16 fields"),
        ("Restaurant", "name, image, rating,\ndeliveryTime, location,\nownerId, categories", CLR_GREEN, "15 fields"),
        ("MenuItem", "restaurantId, name, price,\ncalories, protein, carbs,\nfat, isHealthy", CLR_BLUE, "12 fields"),
        ("Order", "userId, restaurantId,\nitems, status, payment,\ndeliveryAddress", CLR_PURPLE, "20 fields"),
        ("Review", "userId, restaurantId,\nrating, comment,\naiSentiment, images", CLR_CYAN, "14 fields"),
        ("Voucher", "code, type, value,\nminOrder, maxDiscount,\nusageLimit, expiresAt", CLR_ACCENT2, "12 fields"),
        ("GroupOrder", "code, hostId, members,\nitems, status,\nrestaurantId", CLR_RED, "8 fields"),
        ("MealSubscription", "userId, restaurantId,\nplanType, deliveryTime,\ntotalPrice, status", CLR_GREEN, "10 fields"),
        ("Notification", "userId, type, title,\nmessage, read,\ndata", CLR_BLUE, "7 fields"),
        ("CoinTransaction", "userId, amount, coins,\ntype, paymentMethod,\nstatus", CLR_ACCENT, "8 fields"),
        ("PartnerRequest", "userId, restaurantName,\nbusinessType, status,\nownerName, cuisineTypes", CLR_GREEN, "14 fields"),
        ("DriverRequest", "name, email, phone,\nvehicleType, licensePlate,\nstatus", CLR_CYAN, "10 fields"),
    ]
    
    for i, (name, fields, color, count) in enumerate(entities):
        col = i % 4
        row = i // 4
        x = Inches(0.6) + col * Inches(3.1)
        y = Inches(1.4) + row * Inches(1.85)
        w = Inches(2.9)
        h = Inches(1.6)
        
        card = add_shape_with_fill(slide, x, y, w, h, CLR_BG_CARD, 0.05)
        
        # Entity name header
        header = add_shape_with_fill(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.38), color, 0.1)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.8), Inches(0.35),
                     f"📋 {name}", font_size=11, color=CLR_WHITE, bold=True)
        add_text_box(slide, x + Inches(1.6), y + Inches(0.12), Inches(1.1), Inches(0.3),
                     count, font_size=8, color=CLR_WHITE, alignment=PP_ALIGN.RIGHT)
        
        # Fields
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.95),
                     fields, font_size=8, color=CLR_TEXT_DIM, line_spacing=1.5)
    
    # Relationships note
    note_bg = add_shape_with_fill(slide, Inches(0.6), Inches(6.95), Inches(12), Inches(0.45), CLR_BG_CARD2, 0.04)
    add_text_box(slide, Inches(1.0), Inches(6.98), Inches(11.5), Inches(0.35),
                 "🔗 Quan hệ chính: User ←→ Order (1-N)  •  Restaurant ←→ MenuItem (1-N)  •  User ←→ Review (1-N)  •  Order ←→ Review (1-1)  •  Restaurant ←→ Order (1-N)  •  User ←→ Notification (1-N)",
                 font_size=9, color=CLR_TEXT_DIM)


def create_analysis_er_diagram_slide(prs):
    """Slide 12: ER Diagram Visual"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 3, "SƠ ĐỒ QUAN HỆ GIỮA CÁC THỰC THỂ", "Phân tích hệ thống")
    
    # Central entity: User
    center_x, center_y = Inches(5.6), Inches(3.5)
    user_w, user_h = Inches(1.8), Inches(0.7)
    
    user_card = add_shape_with_fill(slide, center_x, center_y, user_w, user_h, CLR_ACCENT, 0.1)
    add_text_box(slide, center_x, center_y + Inches(0.1), user_w, user_h - Inches(0.2),
                 "👤 USER", font_size=14, color=CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Surrounding entities with relationships
    surrounding = [
        ("🛒 Order", Inches(2.0), Inches(1.8), "1 — N", CLR_PURPLE),
        ("🏪 Restaurant", Inches(9.2), Inches(1.8), "1 — N", CLR_GREEN),
        ("⭐ Review", Inches(2.0), Inches(5.0), "1 — N", CLR_CYAN),
        ("🎫 Voucher", Inches(9.2), Inches(5.0), "N — N", CLR_ACCENT2),
        ("🔔 Notification", Inches(0.5), Inches(3.5), "1 — N", CLR_BLUE),
        ("💰 CoinTxn", Inches(10.8), Inches(3.5), "1 — N", CLR_ACCENT),
        ("👥 GroupOrder", Inches(5.6), Inches(1.5), "N — N", CLR_RED),
        ("📅 Subscription", Inches(5.6), Inches(5.5), "1 — N", CLR_GREEN),
    ]
    
    for (label, x, y, rel, color) in surrounding:
        card = add_shape_with_fill(slide, x, y, Inches(1.8), Inches(0.65), CLR_BG_CARD, 0.08)
        # Color top border
        add_gradient_bar(slide, x, y, Inches(1.8), Pt(3), color, color)
        add_text_box(slide, x, y + Inches(0.1), Inches(1.8), Inches(0.25),
                     label, font_size=10, color=CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.35), Inches(1.8), Inches(0.25),
                     rel, font_size=8, color=CLR_TEXT_DIM, alignment=PP_ALIGN.CENTER)
    
    # Restaurant -> MenuItem relationship
    mi_card = add_shape_with_fill(slide, Inches(9.2), Inches(3.2), Inches(1.8), Inches(0.65), CLR_BG_CARD, 0.08)
    add_gradient_bar(slide, Inches(9.2), Inches(3.2), Inches(1.8), Pt(3), CLR_BLUE, CLR_BLUE)
    add_text_box(slide, Inches(9.2), Inches(3.3), Inches(1.8), Inches(0.25),
                 "🍔 MenuItem", font_size=10, color=CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.2), Inches(3.55), Inches(1.8), Inches(0.25),
                 "1 — N", font_size=8, color=CLR_TEXT_DIM, alignment=PP_ALIGN.CENTER)
    
    # Legend
    legend_bg = add_shape_with_fill(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.55), CLR_BG_CARD2, 0.04)
    add_text_box(slide, Inches(1.2), Inches(6.65), Inches(11), Inches(0.4),
                 "📖 Chú thích: 1-N = Một-Nhiều  •  N-N = Nhiều-Nhiều  •  Tổng cộng 16 collections trong MongoDB  •  Sử dụng Mongoose ODM cho schema validation & indexing",
                 font_size=9, color=CLR_TEXT_DIM)


def create_demo_slide(prs):
    """Slide 13: Demo tổng quan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 4, "DEMO SẢN PHẨM", "Trình diễn FoodServe")
    
    features = [
        ("🏠", "Trang chủ", "Hero banner, danh sách nhà hàng,\ntìm kiếm, filter, dark mode,\nchế độ ăn lành mạnh", CLR_ACCENT),
        ("🍔", "Đặt món", "Xem thực đơn, thêm giỏ hàng,\nthanh toán đa phương thức,\nvoucher & mã giảm giá", CLR_GREEN),
        ("📍", "Tracking", "Theo dõi đơn hàng realtime,\nvị trí tài xế trên bản đồ,\ncập nhật trạng thái live", CLR_BLUE),
        ("👥", "Đặt nhóm", "Tạo phòng Split Bill,\nmời bạn bè, đặt chung,\nchia hóa đơn tự động", CLR_PURPLE),
        ("🎮", "Gamification", "Vòng quay may mắn,\nScratch Card, Đố vui,\ntích xu & đổi voucher", CLR_ACCENT2),
        ("🤖", "AI Chatbot", "Hỗ trợ khách hàng 24/7,\ngợi ý món ăn, tra cứu\nđơn hàng (Gemini API)", CLR_CYAN),
        ("🍳", "Quản lý NH", "CRUD menu, xử lý đơn,\nthống kê doanh thu,\nthanh toán phí duy trì", CLR_GREEN),
        ("👑", "Admin Panel", "Dashboard analytics,\nquản lý user/đơn hàng,\nvoucher, cài đặt hệ thống", CLR_RED),
    ]
    
    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 4
        row = i // 4
        x = Inches(0.6) + col * Inches(3.1)
        y = Inches(1.5) + row * Inches(2.75)
        add_icon_card(slide, x, y, Inches(2.85), Inches(2.4), icon, title, desc, color)
    
    # CTA
    cta_bg = add_shape_with_fill(slide, Inches(3.5), Inches(6.85), Inches(6.3), Inches(0.5), CLR_ACCENT, 0.15)
    add_text_box(slide, Inches(3.5), Inches(6.88), Inches(6.3), Inches(0.45),
                 "🎬  LIVE DEMO — TRÌNH DIỄN TRỰC TIẾP TRÊN TRÌNH DUYỆT",
                 font_size=12, color=CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def create_demo_features_slide(prs):
    """Slide 14: Demo - Các tính năng chi tiết"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 4, "TÍNH NĂNG NỔI BẬT", "Trình diễn FoodServe")
    
    # Feature highlight cards
    highlights = [
        ("💳", "THANH TOÁN ĐA KÊNH", 
         "• VNPay (Ngân hàng nội địa)\n• MoMo (Ví điện tử)\n• ZaloPay (QR Code)\n• PayOS (Bank Transfer)\n• Thanh toán bằng Xu\n• Tiền mặt khi nhận hàng",
         CLR_ACCENT),
        ("🗺️", "BẢN ĐỒ & GPS TRACKING",
         "• Leaflet + OpenStreetMap\n• Tracking vị trí tài xế realtime\n• Routing Machine (điều hướng)\n• Tính phí ship theo khoảng cách\n• Geocoding địa chỉ giao hàng\n• Socket.IO live location",
         CLR_BLUE),
        ("🏆", "GAMIFICATION & LOYALTY",
         "• Vòng quay may mắn\n• Scratch Card (Cào thẻ)\n• Quiz đố vui\n• Hệ thống cấp bậc (Rank)\n• Xu thưởng & tích điểm\n• Bảng xếp hạng (Leaderboard)",
         CLR_ACCENT2),
    ]
    
    for i, (icon, title, desc, color) in enumerate(highlights):
        x = Inches(0.6) + i * Inches(4.15)
        y = Inches(1.5)
        
        card = add_shape_with_fill(slide, x, y, Inches(3.85), Inches(3.8), CLR_BG_CARD, 0.04)
        
        # Icon & Title
        icon_bg = add_shape_with_fill(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.6), Inches(0.5), color, 0.1)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.22), Inches(2.6), Inches(0.45),
                     f"{icon}  {title}", font_size=12, color=CLR_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        
        # Description
        add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.3), Inches(2.7),
                     desc, font_size=11, color=CLR_TEXT, line_spacing=1.7)
    
    # Additional features row
    extras = [
        ("🔔", "Thông báo\nRealtime", CLR_GREEN),
        ("🥗", "Chế độ ăn\nLành mạnh", CLR_CYAN),
        ("📅", "Đăng ký\nGói ăn", CLR_PURPLE),
        ("⭐", "Đánh giá\nAI Sentiment", CLR_ACCENT),
        ("❤️", "Nhà hàng\nYêu thích", CLR_RED),
        ("🌙", "Dark Mode\nResponsive", CLR_BLUE),
    ]
    
    for i, (icon, label, color) in enumerate(extras):
        x = Inches(0.6) + i * Inches(2.08)
        y = Inches(5.65)
        card = add_shape_with_fill(slide, x, y, Inches(1.88), Inches(1.2), CLR_BG_CARD2, 0.06)
        add_text_box(slide, x, y + Inches(0.1), Inches(1.88), Inches(0.35),
                     icon, font_size=22, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.55), Inches(1.88), Inches(0.55),
                     label, font_size=9, color=CLR_TEXT, bold=True,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.4)


def create_demo_tech_stats_slide(prs):
    """Slide 15: Thống kê kỹ thuật"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 4, "THỐNG KÊ KỸ THUẬT", "Trình diễn FoodServe")
    
    # Stats cards
    stats = [
        ("25+", "Trang giao diện", "Pages & Components", CLR_ACCENT),
        ("16", "API Routes", "RESTful Endpoints", CLR_BLUE),
        ("16", "Database Models", "MongoDB Collections", CLR_GREEN),
        ("5+", "Cổng thanh toán", "Payment Gateways", CLR_PURPLE),
        ("4", "Vai trò người dùng", "User Roles", CLR_ACCENT2),
        ("50+", "Tính năng", "Features & Functions", CLR_CYAN),
    ]
    
    for i, (number, label, sub, color) in enumerate(stats):
        col = i % 6
        x = Inches(0.6) + col * Inches(2.08)
        y = Inches(1.5)
        
        card = add_shape_with_fill(slide, x, y, Inches(1.88), Inches(1.6), CLR_BG_CARD, 0.06)
        add_text_box(slide, x, y + Inches(0.15), Inches(1.88), Inches(0.6),
                     number, font_size=32, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Segoe UI Black")
        add_text_box(slide, x, y + Inches(0.8), Inches(1.88), Inches(0.3),
                     label, font_size=10, color=CLR_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.1), Inches(1.88), Inches(0.3),
                     sub, font_size=8, color=CLR_TEXT_DIM,
                     alignment=PP_ALIGN.CENTER)
    
    # Technology stack summary
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(6), Inches(0.4),
                 "🔧  TECHNOLOGY STACK TỔNG HỢP", font_size=13, color=CLR_WHITE, bold=True)
    add_decorative_line(slide, Inches(0.8), Inches(3.9), Inches(2), CLR_ACCENT)
    
    stack_items = [
        ("FRONTEND", "React.js 18, Redux Toolkit, React Router v6,\nFramer Motion, Recharts, TailwindCSS,\nSocket.IO Client, Leaflet, Swiper", CLR_BLUE),
        ("BACKEND", "Node.js, Express.js, Mongoose,\nSocket.IO, Nodemailer, bcrypt,\njsonwebtoken, Gemini AI API", CLR_GREEN),
        ("DATABASE", "MongoDB Atlas (NoSQL Cloud),\n16 Collections, Compound Indexes,\nText Search, TTL Indexes", CLR_PURPLE),
        ("DEVOPS", "Vite (Build Tool), Git/GitHub,\nVercel (FE Deploy), Render (BE Deploy),\nPostman (API Testing), ESLint", CLR_ACCENT),
    ]
    
    for i, (label, desc, color) in enumerate(stack_items):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.1)
        y = Inches(4.2) + row * Inches(1.5)
        
        card = add_shape_with_fill(slide, x, y, Inches(5.8), Inches(1.25), CLR_BG_CARD, 0.04)
        
        label_bg = add_shape_with_fill(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.3), Inches(0.35), color, 0.1)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.17), Inches(1.3), Inches(0.3),
                     label, font_size=9, color=CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(1.6), y + Inches(0.15), Inches(3.9), Inches(0.95),
                     desc, font_size=9, color=CLR_TEXT_DIM, line_spacing=1.5)


def create_limitations_slide(prs):
    """Slide 16: Hạn chế & Định hướng"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 5, "HẠN CHẾ & ĐỊNH HƯỚNG PHÁT TRIỂN", "Đánh giá & kế hoạch")
    
    # Left - Limitations
    add_shape_with_fill(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(1.2), Inches(1.65), Inches(4.5), Inches(0.4),
                 "⚠️  HẠN CHẾ HIỆN TẠI", font_size=14, color=CLR_RED, bold=True)
    add_decorative_line(slide, Inches(1.2), Inches(2.15), Inches(1.5), CLR_RED)
    
    limitations = [
        "Chưa có ứng dụng mobile native (iOS/Android)",
        "Chưa tích hợp Push Notification (FCM)",
        "GPS tracking phụ thuộc vào trình duyệt web",
        "Hệ thống AI Chatbot còn giới hạn (chỉ text)",
        "Chưa có hệ thống chat realtime giữa các bên",
        "Chưa tối ưu SEO cho các trang sản phẩm",
        "Chưa có tính năng đánh giá bằng hình ảnh đầy đủ",
        "Chưa có hệ thống báo cáo, xuất PDF tự động",
    ]
    add_multiline_text(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(4),
                       limitations, font_size=11, color=CLR_TEXT, bullet=True, line_spacing=1.7)
    
    # Right - Future
    add_shape_with_fill(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(7.2), Inches(1.65), Inches(4.5), Inches(0.4),
                 "🚀  ĐỊNH HƯỚNG PHÁT TRIỂN", font_size=14, color=CLR_GREEN, bold=True)
    add_decorative_line(slide, Inches(7.2), Inches(2.15), Inches(1.5), CLR_GREEN)
    
    futures = [
        "Phát triển ứng dụng React Native (iOS + Android)",
        "Tích hợp Push Notification qua Firebase Cloud",
        "Nâng cấp AI Chatbot với voice + image recognition",
        "Tích hợp Google Maps / Mapbox cho routing chính xác",
        "Xây dựng hệ thống chat realtime (P2P messaging)",
        "Phát triển hệ thống recommendation AI gợi ý món",
        "Tích hợp nhiều cổng thanh toán quốc tế (Stripe, PayPal)",
        "Triển khai microservices & Docker containerization",
    ]
    add_multiline_text(slide, Inches(7.2), Inches(2.4), Inches(4.8), Inches(4),
                       futures, font_size=11, color=CLR_TEXT, bullet=True, line_spacing=1.7)


def create_conclusion_slide(prs):
    """Slide 17: Kết luận"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    make_slide_header(slide, 6, "KẾT LUẬN", "Tổng kết đồ án")
    
    # Main conclusion card
    main_bg = add_shape_with_fill(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2), CLR_BG_CARD, 0.04)
    add_text_box(slide, Inches(1.3), Inches(1.6), Inches(0.5), Inches(0.5),
                 "✅", font_size=28)
    add_text_box(slide, Inches(1.9), Inches(1.65), Inches(10), Inches(0.4),
                 "Kết quả đạt được", font_size=16, color=CLR_GREEN, bold=True)
    
    conclusion_text = (
        "Đồ án đã hoàn thành xây dựng thành công nền tảng web đặt và giao đồ ăn trực tuyến FoodServe "
        "với đầy đủ các chức năng cốt lõi. Hệ thống phục vụ 4 nhóm người dùng (Khách hàng, Nhà hàng, "
        "Tài xế, Quản trị viên) với giao diện hiện đại, responsive, hỗ trợ dark mode. Tích hợp thành công "
        "các công nghệ tiên tiến: thanh toán đa kênh, tracking realtime, AI Chatbot, gamification, "
        "đặt nhóm Split Bill, và đăng ký gói ăn định kỳ."
    )
    add_text_box(slide, Inches(1.9), Inches(2.15), Inches(10), Inches(1.2),
                 conclusion_text, font_size=12, color=CLR_TEXT, line_spacing=1.7)
    
    # Key achievements
    achievements = [
        ("🏗️", "Fullstack\nhoàn chỉnh", "React.js + Node.js\n+ MongoDB", CLR_BLUE),
        ("⚡", "Realtime\ncommunication", "Socket.IO\nWebSocket", CLR_ACCENT),
        ("💳", "Thanh toán\nđa kênh", "5+ cổng\nthanh toán", CLR_GREEN),
        ("🤖", "AI\nChatbot", "Gemini API\n24/7 hỗ trợ", CLR_PURPLE),
        ("🎮", "Gamification\nSystem", "Vòng quay, cào thẻ\nquiz, rank", CLR_ACCENT2),
        ("📱", "Responsive\nDesign", "Mobile-first\nDark Mode", CLR_CYAN),
    ]
    
    for i, (icon, title, desc, color) in enumerate(achievements):
        x = Inches(0.8) + i * Inches(2.05)
        y = Inches(4.0)
        
        card = add_shape_with_fill(slide, x, y, Inches(1.85), Inches(2.0), CLR_BG_CARD, 0.06)
        add_text_box(slide, x, y + Inches(0.1), Inches(1.85), Inches(0.4),
                     icon, font_size=24, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.6), Inches(1.85), Inches(0.5),
                     title, font_size=10, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.3)
        add_text_box(slide, x, y + Inches(1.2), Inches(1.85), Inches(0.65),
                     desc, font_size=8, color=CLR_TEXT_DIM,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.4)
    
    # Bottom note
    note_bg = add_shape_with_fill(slide, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.55), CLR_BG_CARD2, 0.08)
    add_text_box(slide, Inches(2.5), Inches(6.35), Inches(8.3), Inches(0.45),
                 "💡 Đồ án là sản phẩm thực tế có thể triển khai thương mại, không chỉ dừng ở mức báo cáo học thuật.",
                 font_size=10, color=CLR_ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)


def create_thankyou_slide(prs):
    """Slide 18: Thank you"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_BG_DARK)
    
    # Top accent bar
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CLR_ACCENT, CLR_ACCENT2)
    
    # Decorative circles
    add_decorative_dots(slide, [
        (Inches(0.5), Inches(0.5), Inches(3)),
        (Inches(10), Inches(4.5), Inches(4)),
    ], CLR_ACCENT)
    
    # Big emoji
    add_text_box(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(1.5),
                 "🙏", font_size=80, alignment=PP_ALIGN.CENTER)
    
    # Thank you text
    add_text_box(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(1),
                 "CẢM ƠN HỘI ĐỒNG", font_size=48, color=CLR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI Black")
    
    add_text_box(slide, Inches(0), Inches(3.7), SLIDE_W, Inches(0.6),
                 "ĐÃ LẮNG NGHE BÀI BÁO CÁO", font_size=20, color=CLR_ACCENT,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")
    
    # Decorative line
    add_decorative_line(slide, Inches(5.2), Inches(4.5), Inches(3), CLR_ACCENT)
    
    # Info
    add_text_box(slide, Inches(0), Inches(4.8), SLIDE_W, Inches(0.5),
                 "Sinh viên: Vũ Văn Quyền  •  Đồ án: FoodServe", font_size=14,
                 color=CLR_TEXT_DIM, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0), Inches(5.3), SLIDE_W, Inches(0.5),
                 "Ứng dụng đặt & giao đồ ăn trực tuyến", font_size=12,
                 color=CLR_TEXT_DIM, alignment=PP_ALIGN.CENTER)
    
    # Q&A tag
    qa_bg = add_shape_with_fill(slide, Inches(4.8), Inches(6.0), Inches(3.7), Inches(0.6), CLR_ACCENT, 0.15)
    add_text_box(slide, Inches(4.8), Inches(6.05), Inches(3.7), Inches(0.5),
                 "❓  Q & A  —  HỎI ĐÁP", font_size=14, color=CLR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # Bottom accent bar
    add_gradient_bar(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), CLR_ACCENT, CLR_ACCENT2)


# ============================================================
# MAIN EXECUTION
# ============================================================

def main():
    print("🚀 Đang tạo PowerPoint FoodServe...")
    
    prs = Presentation()
    
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Create all slides
    print("  📋 Slide 1/18: Title / Bìa")
    create_title_slide(prs)
    
    print("  📋 Slide 2/18: Nội dung trình bày")
    create_agenda_slide(prs)
    
    print("  📋 Slide 3/18: Lý do chọn đề tài")
    create_overview_reason_slide(prs)
    
    print("  📋 Slide 4/18: Mục tiêu đề tài")
    create_overview_objectives_slide(prs)
    
    print("  📋 Slide 5/18: Đối tượng & phạm vi")
    create_overview_scope_slide(prs)
    
    print("  📋 Slide 6/18: Phương pháp nghiên cứu")
    create_overview_methodology_slide(prs)
    
    print("  📋 Slide 7/18: Ngôn ngữ lập trình")
    create_theory_languages_slide(prs)
    
    print("  📋 Slide 8/18: Công cụ hỗ trợ")
    create_theory_tools_slide(prs)
    
    print("  📋 Slide 9/18: Môi trường phát triển")
    create_theory_environment_slide(prs)
    
    print("  📋 Slide 10/18: Biểu đồ Use Case")
    create_analysis_usecase_slide(prs)
    
    print("  📋 Slide 11/18: Lược đồ ER (Entities)")
    create_analysis_er_slide(prs)
    
    print("  📋 Slide 12/18: Sơ đồ quan hệ ER")
    create_analysis_er_diagram_slide(prs)
    
    print("  📋 Slide 13/18: Demo tổng quan")
    create_demo_slide(prs)
    
    print("  📋 Slide 14/18: Tính năng nổi bật")
    create_demo_features_slide(prs)
    
    print("  📋 Slide 15/18: Thống kê kỹ thuật")
    create_demo_tech_stats_slide(prs)
    
    print("  📋 Slide 16/18: Hạn chế & Định hướng")
    create_limitations_slide(prs)
    
    print("  📋 Slide 17/18: Kết luận")
    create_conclusion_slide(prs)
    
    print("  📋 Slide 18/18: Cảm ơn & Q&A")
    create_thankyou_slide(prs)
    
    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'FoodServe_Presentation.pptx')
    prs.save(output_path)
    print(f"\n✅ Đã tạo thành công: {output_path}")
    print(f"📊 Tổng cộng: 18 slides")
    print(f"🎨 Theme: Dark Premium (Cam & Đen)")
    print(f"📐 Kích thước: Widescreen 16:9")


if __name__ == '__main__':
    main()
